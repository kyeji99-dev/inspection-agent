"""
AI Vision 기반 결함 검사 보조 에이전트
- 공개 산업 결함 데이터셋(MVTec AD 등) 기반 데모용
- Claude Vision으로 이미지 분석 → 웹 검색으로 원인 조사 → PPT 보고서 자동 생성
"""

import base64
import io
import json
import re
from datetime import datetime

import anthropic
import streamlit as st
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# =============================================================================
# 페이지 설정
# =============================================================================
st.set_page_config(
    page_title="Vision Inspection Agent",
    page_icon="🔍",
    layout="wide",
)

st.markdown(
    """
    <style>
    .block-container { padding-top: 2rem; padding-bottom: 2rem; }
    .stButton > button { font-weight: 600; }
    div[data-testid="stMetricValue"] { font-size: 1.8rem; }
    </style>
    """,
    unsafe_allow_html=True,
)

st.title("🔍 Vision Inspection Agent")
st.caption(
    "AI 비전 기반 결함 검사 · 자동 원인 분석 · 임원 보고서 생성   |   Powered by Claude"
)

# =============================================================================
# 사이드바
# =============================================================================
with st.sidebar:
    st.header("⚙️ 설정")
    api_key = st.text_input(
        "Anthropic API Key",
        type="password",
        help="https://console.anthropic.com 에서 발급",
    )
    model = st.selectbox(
        "Model",
        ["claude-opus-4-7", "claude-sonnet-4-6"],
        help="Opus = 더 정확 / Sonnet = 더 빠름",
    )

    st.divider()
    st.subheader("🏭 검사 대상")
    product_type = st.text_input(
        "제품/소재 (예: PCB, 가죽, 직물, OLED 패널)",
        value="PCB",
    )
    inspection_focus = st.text_input(
        "중점 검사 항목",
        value="표면 결함, 오염, 스크래치, 균열",
    )

    st.divider()
    st.subheader("📂 데모 데이터")
    st.markdown(
        """
        본 데모는 **공개 데이터셋**만 사용합니다:
        - [MVTec AD](https://www.mvtec.com/company/research/datasets/mvtec-ad)
        - [KolektorSDD](https://www.vicos.si/resources/kolektorsdd/)
        - Hugging Face: `Voxel51/mvtec-ad`
        """
    )

# =============================================================================
# 디자인 토큰
# =============================================================================
NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
MUTED = RGBColor(0x66, 0x66, 0x80)

SEV_HIGH = RGBColor(0xC1, 0x21, 0x1B)
SEV_MED = RGBColor(0xE6, 0x9F, 0x00)
SEV_LOW = RGBColor(0x4E, 0x7A, 0x52)
SEV_NONE = RGBColor(0x4E, 0x7A, 0x52)

SEVERITY_COLORS_HEX = {
    "HIGH": "#C1211B",
    "MEDIUM": "#E69F00",
    "LOW": "#4E7A52",
    "NONE": "#4E7A52",
}


# =============================================================================
# 핵심 1: 이미지 분석 (Claude Vision)
# =============================================================================
def image_to_base64(img_bytes: bytes) -> str:
    return base64.standard_b64encode(img_bytes).decode("utf-8")


def detect_media_type(filename: str) -> str:
    f = filename.lower()
    if f.endswith(".png"):
        return "image/png"
    if f.endswith((".jpg", ".jpeg")):
        return "image/jpeg"
    if f.endswith(".webp"):
        return "image/webp"
    if f.endswith(".gif"):
        return "image/gif"
    return "image/jpeg"


VISION_PROMPT_TEMPLATE = """당신은 **{product_type}** 제조 품질 검사 전문가입니다.
주어진 이미지를 산업 검사 기준으로 정밀 분석하세요.

중점 검사 항목: {inspection_focus}

## 출력 형식 (반드시 ```json 코드 블록```으로 감싸서 출력)

```json
{{
  "has_defect": true/false,
  "overall_assessment": "검사 결과 한 줄 요약 (한국어)",
  "overall_severity": "HIGH | MEDIUM | LOW | NONE",
  "confidence": 0.0~1.0,
  "defects": [
    {{
      "type": "결함 유형 (한국어, 예: 스크래치, 오염, 균열, 변색, 핀홀, 박리)",
      "type_en": "Defect type in English",
      "location_desc": "결함 위치 설명 (예: 우상단, 중앙부, 좌측 모서리)",
      "bbox_normalized": [x1, y1, x2, y2],
      "size_estimate": "결함 크기 추정 (정성적, 예: 미세, 중간, 큰)",
      "description": "결함 상세 묘사 (1~2문장)",
      "severity": "HIGH | MEDIUM | LOW",
      "possible_causes": ["원인 추정 1", "원인 추정 2", "원인 추정 3"]
    }}
  ],
  "inspection_notes": "검사관 톤의 종합 소견 (3~4문장, 임원 보고용)",
  "recommended_actions": [
    "권장 조치 1 (구체적, 실행 가능)",
    "권장 조치 2",
    "권장 조치 3"
  ]
}}
```

## 분석 가이드
- bbox_normalized는 0~1 정규화 좌표 [좌상x, 좌상y, 우하x, 우하y]
  (이미지 전체가 (0,0)~(1,1), Y는 위에서 아래로 증가)
- 결함이 여러 개면 모두 나열
- 결함이 없으면 has_defect=false, defects=[]
- 추측은 신중하게, 불확실하면 confidence를 낮게
- 한국어 임원 보고 톤 사용
"""


def analyze_image(client, model, image_bytes, media_type, product_type, focus):
    prompt = VISION_PROMPT_TEMPLATE.format(
        product_type=product_type,
        inspection_focus=focus,
    )

    resp = client.messages.create(
        model=model,
        max_tokens=4000,
        messages=[
            {
                "role": "user",
                "content": [
                    {
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": media_type,
                            "data": image_to_base64(image_bytes),
                        },
                    },
                    {"type": "text", "text": prompt},
                ],
            }
        ],
    )

    text = "".join(b.text for b in resp.content if getattr(b, "type", None) == "text")
    m = re.search(r"```json\s*(.*?)\s*```", text, re.DOTALL)
    if not m:
        raise ValueError(f"JSON 형식을 찾을 수 없습니다:\n{text}")
    return json.loads(m.group(1))


# =============================================================================
# 핵심 2: 웹 리서치 (근본 원인 + 유사 사례)
# =============================================================================
RESEARCH_PROMPT_TEMPLATE = """다음 제조 결함에 대해 웹 검색으로 산업 표준 정보를 조사해주세요.

## 분석 대상
- 제품/소재: {product_type}
- 발견된 결함 유형: {defect_types}
- 종합 심각도: {severity}

## 조사 항목 (웹 검색 활용)
1. 이 유형의 결함의 **공정상 일반적인 근본 원인** (구체적, 데이터 기반)
2. **산업 표준 검사 방법** 또는 관련 ISO/IEC 표준
3. **유사 사례** (논문, 산업 케이스 스터디, 기술 블로그)
4. **예방·개선 모범 사례** (공정 설계, 검사 강화 방안)

## 출력 형식 (반드시 ```json 코드 블록```으로 감싸서 출력)

```json
{{
  "root_cause_analysis": [
    {{"category": "공정 요인", "details": "1~2문장 한국어"}},
    {{"category": "설비 요인", "details": "..."}},
    {{"category": "재료 요인", "details": "..."}},
    {{"category": "환경 요인", "details": "..."}}
  ],
  "industry_standards": [
    "관련 표준/규격 (간단한 설명 포함)"
  ],
  "similar_cases": [
    "유사 사례 또는 논문 요약 (1~2문장)"
  ],
  "best_practices": [
    "모범 사례 1",
    "모범 사례 2",
    "모범 사례 3"
  ],
  "improvement_recommendations": [
    "개선 권고 1 (구체적, 우선순위 높음)",
    "개선 권고 2",
    "개선 권고 3"
  ]
}}
```

모든 텍스트는 한국어로, 임원 보고 톤(간결, 데이터 중심)으로 작성하세요.
"""


def research_root_cause(client, model, product_type, defects, severity):
    defect_types = ", ".join(
        sorted({d.get("type", "") for d in defects if d.get("type")})
    ) or "특이사항 없음"

    prompt = RESEARCH_PROMPT_TEMPLATE.format(
        product_type=product_type,
        defect_types=defect_types,
        severity=severity,
    )

    resp = client.messages.create(
        model=model,
        max_tokens=8000,
        tools=[
            {
                "type": "web_search_20250305",
                "name": "web_search",
                "max_uses": 6,
            }
        ],
        messages=[{"role": "user", "content": prompt}],
    )

    text = ""
    queries = []
    sources = []
    for block in resp.content:
        btype = getattr(block, "type", None)
        if btype == "text":
            text += block.text + "\n"
        elif btype == "server_tool_use" and getattr(block, "name", "") == "web_search":
            q = (block.input or {}).get("query", "")
            if q:
                queries.append(q)
        elif btype == "web_search_tool_result":
            for item in getattr(block, "content", []) or []:
                url = (
                    getattr(item, "url", None)
                    if not isinstance(item, dict)
                    else item.get("url")
                )
                if url and url not in sources:
                    sources.append(url)

    m = re.search(r"```json\s*(.*?)\s*```", text, re.DOTALL)
    if not m:
        return {"_raw": text, "_queries": queries, "_sources": sources}

    data = json.loads(m.group(1))
    data["_queries"] = queries
    data["_sources"] = sources
    return data


# =============================================================================
# 결함 표시 이미지 생성
# =============================================================================
def draw_defect_overlay(img_bytes: bytes, defects: list) -> bytes:
    img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
    draw = ImageDraw.Draw(img)
    W, H = img.size

    palette = {
        "HIGH": (193, 33, 27),
        "MEDIUM": (230, 159, 0),
        "LOW": (78, 122, 82),
    }

    for i, d in enumerate(defects):
        bbox = d.get("bbox_normalized")
        if not bbox or len(bbox) != 4:
            continue
        sev = d.get("severity", "MEDIUM")
        color = palette.get(sev, (193, 33, 27))

        x1 = int(bbox[0] * W)
        y1 = int(bbox[1] * H)
        x2 = int(bbox[2] * W)
        y2 = int(bbox[3] * H)

        line_w = max(3, min(W, H) // 200)
        draw.rectangle([x1, y1, x2, y2], outline=color, width=line_w)

        label = f"#{i+1} {d.get('type_en', d.get('type', 'Defect'))}"
        try:
            tbbox = draw.textbbox((0, 0), label)
            tw, th = tbbox[2] - tbbox[0], tbbox[3] - tbbox[1]
        except AttributeError:
            tw, th = 8 * len(label), 14

        pad = 4
        ly1 = max(0, y1 - th - pad * 2)
        draw.rectangle(
            [x1, ly1, x1 + tw + pad * 2, ly1 + th + pad * 2],
            fill=color,
        )
        draw.text((x1 + pad, ly1 + pad), label, fill=(255, 255, 255))

    out = io.BytesIO()
    img.save(out, format="PNG")
    out.seek(0)
    return out.getvalue()


# =============================================================================
# PPT 보고서 생성
# =============================================================================
def _fill_bg(slide, color, prs):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def _add_text(
    slide, left, top, width, height, text,
    size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.text = text
    for p in tf.paragraphs:
        p.alignment = align
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = font
    return box


def severity_to_color(sev: str) -> RGBColor:
    return {
        "HIGH": SEV_HIGH,
        "MEDIUM": SEV_MED,
        "LOW": SEV_LOW,
        "NONE": SEV_NONE,
    }.get(sev, SEV_MED)


def generate_inspection_pptx(
    original_bytes: bytes,
    annotated_bytes: bytes,
    analysis: dict,
    research: dict,
    product_type: str,
) -> io.BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    today = datetime.now().strftime("%Y.%m.%d")
    severity = analysis.get("overall_severity", "NONE")

    # ───── 슬라이드 1: 표지 (다크) ─────
    s = prs.slides.add_slide(blank)
    _fill_bg(s, NAVY, prs)

    _add_text(
        s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.4),
        "QUALITY INSPECTION REPORT",
        size=11, bold=True, color=ICE,
    )
    _add_text(
        s, Inches(0.8), Inches(2.7), Inches(11.7), Inches(2.5),
        f"{product_type} 검사 결과 보고",
        size=44, bold=True, color=WHITE,
    )
    _add_text(
        s, Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.5),
        analysis.get("overall_assessment", ""),
        size=18, color=ICE,
    )
    _add_text(
        s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.4),
        f"Inspection Date  {today}   ·   AI Vision Inspection Agent",
        size=10, color=ICE,
    )

    # ───── 슬라이드 2: 검사 요약 + 원본/주석 이미지 ─────
    s = prs.slides.add_slide(blank)
    _fill_bg(s, WHITE, prs)

    _add_text(
        s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.3),
        "01  ·  INSPECTION SUMMARY",
        size=10, bold=True, color=NAVY,
    )
    _add_text(
        s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.6),
        "검사 요약",
        size=28, bold=True, color=DARK,
    )

    sev_color = severity_to_color(severity)
    _add_text(
        s, Inches(0.7), Inches(1.95), Inches(3.5), Inches(0.35),
        "OVERALL SEVERITY",
        size=10, bold=True, color=MUTED,
    )
    _add_text(
        s, Inches(0.7), Inches(2.3), Inches(3.5), Inches(0.9),
        severity,
        size=44, bold=True, color=sev_color,
    )

    n_defects = len(analysis.get("defects", []))
    _add_text(
        s, Inches(0.7), Inches(3.5), Inches(3.5), Inches(0.35),
        "DEFECTS DETECTED",
        size=10, bold=True, color=MUTED,
    )
    _add_text(
        s, Inches(0.7), Inches(3.85), Inches(3.5), Inches(0.9),
        f"{n_defects}",
        size=44, bold=True, color=NAVY,
    )

    conf = analysis.get("confidence", 0)
    _add_text(
        s, Inches(0.7), Inches(5.05), Inches(3.5), Inches(0.35),
        "CONFIDENCE",
        size=10, bold=True, color=MUTED,
    )
    _add_text(
        s, Inches(0.7), Inches(5.4), Inches(3.5), Inches(0.9),
        f"{conf*100:.0f}%",
        size=44, bold=True, color=NAVY,
    )

    img_bytes_for_slide = annotated_bytes if annotated_bytes else original_bytes
    img_stream = io.BytesIO(img_bytes_for_slide)
    try:
        s.shapes.add_picture(
            img_stream, Inches(5.0), Inches(2.0),
            width=Inches(7.8), height=Inches(5.0),
        )
        _add_text(
            s, Inches(5.0), Inches(7.05), Inches(7.8), Inches(0.3),
            "▲ AI 검출 결함 영역 표시" if annotated_bytes else "▲ 검사 이미지",
            size=10, color=MUTED,
        )
    except Exception:
        pass

    # ───── 슬라이드 3: 결함 상세 ─────
    s = prs.slides.add_slide(blank)
    _fill_bg(s, WHITE, prs)

    _add_text(
        s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.3),
        "02  ·  DEFECT DETAILS",
        size=10, bold=True, color=NAVY,
    )
    _add_text(
        s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.6),
        "결함 상세 분석",
        size=28, bold=True, color=DARK,
    )

    defects = analysis.get("defects", [])[:4]
    if defects:
        for i, d in enumerate(defects):
            row, col = divmod(i, 2)
            left = Inches(0.7 + col * 6.1)
            top = Inches(2.0 + row * 2.55)

            sev = d.get("severity", "MEDIUM")
            sev_c = severity_to_color(sev)

            bar = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top,
                Inches(0.1), Inches(2.3),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = sev_c
            bar.line.fill.background()

            _add_text(
                s, left + Inches(0.3), top - Inches(0.05),
                Inches(5.4), Inches(0.5),
                f"#{i+1}  {d.get('type', '')}",
                size=18, bold=True, color=DARK,
            )

            _add_text(
                s, left + Inches(0.3), top + Inches(0.45),
                Inches(2.5), Inches(0.3),
                f"SEVERITY  {sev}",
                size=10, bold=True, color=sev_c,
            )

            _add_text(
                s, left + Inches(0.3), top + Inches(0.8),
                Inches(5.4), Inches(0.3),
                f"위치: {d.get('location_desc', '-')}",
                size=11, color=DARK,
            )

            _add_text(
                s, left + Inches(0.3), top + Inches(1.1),
                Inches(5.4), Inches(0.3),
                f"크기: {d.get('size_estimate', '-')}",
                size=11, color=DARK,
            )

            desc = d.get("description", "")
            if len(desc) > 130:
                desc = desc[:130] + "..."
            _add_text(
                s, left + Inches(0.3), top + Inches(1.45),
                Inches(5.4), Inches(0.8),
                desc, size=11, color=DARK,
            )
    else:
        _add_text(
            s, Inches(0.7), Inches(3.5), Inches(12), Inches(1),
            "특이 결함이 검출되지 않았습니다",
            size=24, bold=True, color=SEV_LOW, align=PP_ALIGN.CENTER,
        )

    # ───── 슬라이드 4: 근본 원인 분석 ─────
    s = prs.slides.add_slide(blank)
    _fill_bg(s, WHITE, prs)

    _add_text(
        s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.3),
        "03  ·  ROOT CAUSE ANALYSIS",
        size=10, bold=True, color=NAVY,
    )
    _add_text(
        s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.6),
        "근본 원인 분석",
        size=28, bold=True, color=DARK,
    )
    _add_text(
        s, Inches(0.7), Inches(1.7), Inches(12), Inches(0.4),
        "웹 검색을 통해 자동 수집된 산업 표준 기반 분석",
        size=11, color=MUTED,
    )

    causes = research.get("root_cause_analysis", [])[:4]
    for i, c in enumerate(causes):
        row, col = divmod(i, 2)
        left = Inches(0.7 + col * 6.1)
        top = Inches(2.4 + row * 2.3)

        num = s.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.5), Inches(0.5))
        num.fill.solid()
        num.fill.fore_color.rgb = NAVY
        num.line.fill.background()
        tf = num.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        tf.text = str(i + 1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for r in p.runs:
            r.font.size = Pt(16)
            r.font.bold = True
            r.font.color.rgb = WHITE

        _add_text(
            s, left + Inches(0.7), top - Inches(0.02),
            Inches(5.0), Inches(0.45),
            c.get("category", ""),
            size=15, bold=True, color=NAVY,
        )

        details = c.get("details", "")
        if len(details) > 180:
            details = details[:180] + "..."
        _add_text(
            s, left + Inches(0.7), top + Inches(0.5),
            Inches(5.0), Inches(1.7),
            details, size=11, color=DARK,
        )

    # ───── 슬라이드 5: 권고 사항 ─────
    s = prs.slides.add_slide(blank)
    _fill_bg(s, WHITE, prs)

    _add_text(
        s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.3),
        "04  ·  RECOMMENDATIONS",
        size=10, bold=True, color=NAVY,
    )
    _add_text(
        s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.6),
        "권고 조치",
        size=28, bold=True, color=DARK,
    )

    actions = list(analysis.get("recommended_actions", []))
    actions += list(research.get("improvement_recommendations", []))
    actions = actions[:5]

    for i, a in enumerate(actions):
        top = Inches(1.95 + i * 1.0)
        _add_text(
            s, Inches(0.7), top - Inches(0.05),
            Inches(0.9), Inches(0.9),
            f"{i+1:02d}", size=32, bold=True, color=ICE,
        )
        if len(a) > 160:
            a = a[:160] + "..."
        _add_text(
            s, Inches(1.7), top + Inches(0.1),
            Inches(11.0), Inches(0.9),
            a, size=14, color=DARK,
        )

    # ───── 슬라이드 6: 참고 자료 (다크) ─────
    s = prs.slides.add_slide(blank)
    _fill_bg(s, NAVY, prs)

    _add_text(
        s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.3),
        "05  ·  REFERENCES",
        size=10, bold=True, color=ICE,
    )
    _add_text(
        s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.6),
        "참고 자료",
        size=28, bold=True, color=WHITE,
    )
    _add_text(
        s, Inches(0.7), Inches(1.85), Inches(12), Inches(0.4),
        "산업 표준 · 본 검사 보고서의 근거 자료",
        size=12, color=ICE,
    )

    _add_text(
        s, Inches(0.7), Inches(2.6), Inches(5.8), Inches(0.35),
        "INDUSTRY STANDARDS",
        size=11, bold=True, color=ICE,
    )
    standards = research.get("industry_standards", [])[:6]
    standards_text = "\n".join([f"·  {st}" for st in standards]) or "·  (해당 사항 없음)"
    _add_text(
        s, Inches(0.7), Inches(3.0), Inches(5.8), Inches(3.5),
        standards_text, size=10, color=ICE,
    )

    _add_text(
        s, Inches(6.9), Inches(2.6), Inches(5.8), Inches(0.35),
        "WEB SOURCES",
        size=11, bold=True, color=ICE,
    )
    sources = research.get("_sources", [])[:8]
    src_lines = []
    for i, src in enumerate(sources, 1):
        s_disp = src if len(src) <= 70 else src[:67] + "..."
        src_lines.append(f"[{i:02d}]  {s_disp}")
    src_text = "\n".join(src_lines) or "(검색 결과 없음)"
    _add_text(
        s, Inches(6.9), Inches(3.0), Inches(5.8), Inches(3.5),
        src_text, size=9, color=ICE,
    )

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# =============================================================================
# 메인 UI
# =============================================================================
st.subheader("📤 검사 이미지 업로드")

uploaded = st.file_uploader(
    "공개 데이터셋 이미지 (MVTec AD 등)를 업로드하세요",
    type=["png", "jpg", "jpeg", "webp"],
)

col_btn1, _ = st.columns([1, 4])
with col_btn1:
    run = st.button("🚀 검사 시작", type="primary", use_container_width=True)

if uploaded:
    preview = Image.open(uploaded)
    st.image(preview, caption=f"업로드: {uploaded.name}", width=400)

if run:
    if not api_key:
        st.error("⚠️ 사이드바에서 Anthropic API 키를 입력해주세요")
        st.stop()
    if not uploaded:
        st.error("⚠️ 검사할 이미지를 업로드해주세요")
        st.stop()

    uploaded.seek(0)
    image_bytes = uploaded.read()
    media_type = detect_media_type(uploaded.name)
    client = anthropic.Anthropic(api_key=api_key)

    # ── 1단계: 비전 분석 ─────────────────────────────
    with st.status("🔍 1단계 · Claude Vision으로 결함 분석 중...", expanded=True) as stt:
        st.write("이미지를 정밀 분석하여 결함을 검출합니다.")
        try:
            analysis = analyze_image(
                client, model, image_bytes, media_type,
                product_type, inspection_focus,
            )
        except Exception as e:
            st.error(f"비전 분석 실패: {e}")
            st.stop()
        st.write(
            f"✅ 결함 **{len(analysis.get('defects', []))}건** 검출 · "
            f"종합 심각도 **{analysis.get('overall_severity', 'NONE')}**"
        )
        stt.update(label="✅ 1단계 완료 · 비전 분석", state="complete")

    defects = analysis.get("defects", [])
    annotated_bytes = None
    if defects:
        try:
            annotated_bytes = draw_defect_overlay(image_bytes, defects)
        except Exception as e:
            st.warning(f"오버레이 생성 실패 (계속 진행): {e}")

    # ── 2단계: 웹 리서치 ─────────────────────────────
    with st.status("🌐 2단계 · 근본 원인 자동 조사 중...", expanded=True) as stt:
        st.write("결함 유형 기반으로 웹 검색을 수행해 산업 표준·원인 정보를 수집합니다.")
        try:
            research = research_root_cause(
                client, model, product_type,
                analysis.get("defects", []),
                analysis.get("overall_severity", "NONE"),
            )
        except Exception as e:
            st.error(f"리서치 실패: {e}")
            research = {"_sources": [], "_queries": []}
        st.write(
            f"🔎 검색 **{len(research.get('_queries', []))}회** · "
            f"출처 **{len(research.get('_sources', []))}개** 수집"
        )
        with st.expander("수행된 검색 쿼리"):
            for q in research.get("_queries", []):
                st.code(q, language=None)
        stt.update(label="✅ 2단계 완료 · 원인 조사", state="complete")

    # ── 3단계: PPT 생성 ─────────────────────────────
    with st.status("📊 3단계 · 임원 보고용 PPT 생성 중...", expanded=False) as stt:
        ppt_buf = generate_inspection_pptx(
            image_bytes, annotated_bytes, analysis, research, product_type,
        )
        stt.update(label="✅ 3단계 완료 · 보고서 생성", state="complete")

    st.success("🎉 검사 완료")

    fname_safe = re.sub(r"[^\w가-힣]+", "_", uploaded.name.rsplit(".", 1)[0])
    st.download_button(
        "📥 검사 보고서 (.pptx) 다운로드",
        data=ppt_buf,
        file_name=f"검사보고서_{product_type}_{fname_safe}.pptx",
        mime=(
            "application/vnd.openxmlformats-officedocument."
            "presentationml.presentation"
        ),
        type="primary",
        use_container_width=True,
    )

    st.divider()

    sev = analysis.get("overall_severity", "NONE")
    c1, c2, c3, c4 = st.columns(4)
    c1.metric("종합 심각도", sev)
    c2.metric("검출 결함", f"{len(defects)}건")
    c3.metric("신뢰도", f"{analysis.get('confidence', 0)*100:.0f}%")
    c4.metric("검사 일자", datetime.now().strftime("%Y-%m-%d"))

    st.markdown("### 종합 평가")
    st.info(analysis.get("overall_assessment", ""))

    st.markdown("### 검사 이미지")
    img_col1, img_col2 = st.columns(2)
    with img_col1:
        st.image(image_bytes, caption="원본 이미지", use_column_width=True)
    with img_col2:
        if annotated_bytes:
            st.image(
                annotated_bytes,
                caption="AI 검출 결함 영역",
                use_column_width=True,
            )
        else:
            st.info("결함이 검출되지 않았습니다.")

    if defects:
        st.markdown("### 결함 상세")
        for i, d in enumerate(defects, 1):
            with st.expander(
                f"#{i}  {d.get('type', '')}  ·  Severity: {d.get('severity', '-')}"
            ):
                col_a, col_b = st.columns(2)
                col_a.write(f"**위치:** {d.get('location_desc', '-')}")
                col_b.write(f"**크기:** {d.get('size_estimate', '-')}")
                st.write(f"**설명:** {d.get('description', '-')}")
                st.write("**추정 원인:**")
                for c in d.get("possible_causes", []):
                    st.write(f"  - {c}")

    st.markdown("### 검사관 소견")
    st.write(analysis.get("inspection_notes", ""))

    if research.get("root_cause_analysis"):
        st.markdown("### 근본 원인 분석 (웹 리서치)")
        for c in research["root_cause_analysis"]:
            st.write(f"**{c.get('category', '')}** — {c.get('details', '')}")

    st.markdown("### 권고 조치")
    actions = list(analysis.get("recommended_actions", []))
    actions += list(research.get("improvement_recommendations", []))
    for i, a in enumerate(actions[:5], 1):
        st.write(f"**{i}.** {a}")

    with st.expander("📚 산업 표준 / 모범 사례"):
        st.markdown("**관련 표준**")
        for s in research.get("industry_standards", []):
            st.write(f"- {s}")
        st.markdown("**모범 사례**")
        for s in research.get("best_practices", []):
            st.write(f"- {s}")
        st.markdown("**유사 사례**")
        for s in research.get("similar_cases", []):
            st.write(f"- {s}")

    with st.expander("🔗 웹 출처"):
        for i, src in enumerate(research.get("_sources", []), 1):
            st.write(f"[{i}] {src}")

    with st.expander("🧪 Raw JSON (디버깅용)"):
        st.json({"analysis": analysis, "research": research})
