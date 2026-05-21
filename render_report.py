"""
Render an inspection report PPT from analysis + research JSON.

No API keys or network access required. Invoked by the /inspect slash
command after Claude Code has produced the two JSON files.

Usage:
    python render_report.py <image> <analysis.json> <research.json> \
        --product "PCB" --out report.pptx --annotated-out annotated.png
"""

import argparse
import io
import json
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# =============================================================================
# 디자인 토큰
# =============================================================================
NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
MUTED = RGBColor(0x66, 0x66, 0x80)

SEV_HIGH = RGBColor(0xC1, 0x21, 0x1B)
SEV_MED = RGBColor(0xE6, 0x9F, 0x00)
SEV_LOW = RGBColor(0x4E, 0x7A, 0x52)
SEV_NONE = RGBColor(0x4E, 0x7A, 0x52)


# =============================================================================
# 결함 표시 이미지 생성
# =============================================================================
def draw_defect_overlay(img_bytes: bytes, defects: list) -> bytes:
    img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
    draw = ImageDraw.Draw(img)
    W, H = img.size

    palette = {
        "HIGH": (193, 33, 27),
        "MEDIUM": (230, 159, 0),
        "LOW": (78, 122, 82),
    }

    for i, d in enumerate(defects):
        bbox = d.get("bbox_normalized")
        if not bbox or len(bbox) != 4:
            continue
        sev = d.get("severity", "MEDIUM")
        color = palette.get(sev, (193, 33, 27))

        x1 = int(bbox[0] * W)
        y1 = int(bbox[1] * H)
        x2 = int(bbox[2] * W)
        y2 = int(bbox[3] * H)

        line_w = max(3, min(W, H) // 200)
        draw.rectangle([x1, y1, x2, y2], outline=color, width=line_w)

        label = f"#{i+1} {d.get('type_en', d.get('type', 'Defect'))}"
        try:
            tbbox = draw.textbbox((0, 0), label)
            tw, th = tbbox[2] - tbbox[0], tbbox[3] - tbbox[1]
        except AttributeError:
            tw, th = 8 * len(label), 14

        pad = 4
        ly1 = max(0, y1 - th - pad * 2)
        draw.rectangle(
            [x1, ly1, x1 + tw + pad * 2, ly1 + th + pad * 2],
            fill=color,
        )
        draw.text((x1 + pad, ly1 + pad), label, fill=(255, 255, 255))

    out = io.BytesIO()
    img.save(out, format="PNG")
    out.seek(0)
    return out.getvalue()


# =============================================================================
# PPT 보고서 생성
# =============================================================================
def _fill_bg(slide, color, prs):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def _add_text(
    slide, left, top, width, height, text,
    size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.text = text
    for p in tf.paragraphs:
        p.alignment = align
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = font
    return box


def severity_to_color(sev: str) -> RGBColor:
    return {
        "HIGH": SEV_HIGH,
        "MEDIUM": SEV_MED,
        "LOW": SEV_LOW,
        "NONE": SEV_NONE,
    }.get(sev, SEV_MED)


def generate_inspection_pptx(
    original_bytes: bytes,
    annotated_bytes: bytes,
    analysis: dict,
    research: dict,
    product_type: str,
) -> io.BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    from datetime import datetime
    today = datetime.now().strftime("%Y.%m.%d")
    severity = analysis.get("overall_severity", "NONE")

    # ───── 슬라이드 1: 표지 (다크) ─────
    s = prs.slides.add_slide(blank)
    _fill_bg(s, NAVY, prs)

    _add_text(
        s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.4),
        "QUALITY INSPECTION REPORT",
        size=11, bold=True, color=ICE,
    )
    _add_text(
        s, Inches(0.8), Inches(2.7), Inches(11.7), Inches(2.5),
        f"{product_type} 검사 결과 보고",
        size=44, bold=True, color=WHITE,
    )
    _add_text(
        s, Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.5),
        analysis.get("overall_assessment", ""),
        size=18, color=ICE,
    )
    _add_text(
        s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.4),
        f"Inspection Date  {today}   ·   AI Vision Inspection Agent",
        size=10, color=ICE,
    )

    # ───── 슬라이드 2: 검사 요약 + 원본/주석 이미지 ─────
    s = prs.slides.add_slide(blank)
    _fill_bg(s, WHITE, prs)

    _add_text(
        s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.3),
        "01  ·  INSPECTION SUMMARY",
        size=10, bold=True, color=NAVY,
    )
    _add_text(
        s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.6),
        "검사 요약",
        size=28, bold=True, color=DARK,
    )

    sev_color = severity_to_color(severity)
    _add_text(
        s, Inches(0.7), Inches(1.95), Inches(3.5), Inches(0.35),
        "OVERALL SEVERITY",
        size=10, bold=True, color=MUTED,
    )
    _add_text(
        s, Inches(0.7), Inches(2.3), Inches(3.5), Inches(0.9),
        severity,
        size=44, bold=True, color=sev_color,
    )

    n_defects = len(analysis.get("defects", []))
    _add_text(
        s, Inches(0.7), Inches(3.5), Inches(3.5), Inches(0.35),
        "DEFECTS DETECTED",
        size=10, bold=True, color=MUTED,
    )
    _add_text(
        s, Inches(0.7), Inches(3.85), Inches(3.5), Inches(0.9),
        f"{n_defects}",
        size=44, bold=True, color=NAVY,
    )

    conf = analysis.get("confidence", 0)
    _add_text(
        s, Inches(0.7), Inches(5.05), Inches(3.5), Inches(0.35),
        "CONFIDENCE",
        size=10, bold=True, color=MUTED,
    )
    _add_text(
        s, Inches(0.7), Inches(5.4), Inches(3.5), Inches(0.9),
        f"{conf*100:.0f}%",
        size=44, bold=True, color=NAVY,
    )

    img_bytes_for_slide = annotated_bytes if annotated_bytes else original_bytes
    img_stream = io.BytesIO(img_bytes_for_slide)
    try:
        s.shapes.add_picture(
            img_stream, Inches(5.0), Inches(2.0),
            width=Inches(7.8), height=Inches(5.0),
        )
        _add_text(
            s, Inches(5.0), Inches(7.05), Inches(7.8), Inches(0.3),
            "▲ AI 검출 결함 영역 표시" if annotated_bytes else "▲ 검사 이미지",
            size=10, color=MUTED,
        )
    except Exception:
        pass

    # ───── 슬라이드 3: 결함 상세 ─────
    s = prs.slides.add_slide(blank)
    _fill_bg(s, WHITE, prs)

    _add_text(
        s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.3),
        "02  ·  DEFECT DETAILS",
        size=10, bold=True, color=NAVY,
    )
    _add_text(
        s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.6),
        "결함 상세 분석",
        size=28, bold=True, color=DARK,
    )

    defects = analysis.get("defects", [])[:4]
    if defects:
        for i, d in enumerate(defects):
            row, col = divmod(i, 2)
            left = Inches(0.7 + col * 6.1)
            top = Inches(2.0 + row * 2.55)

            sev = d.get("severity", "MEDIUM")
            sev_c = severity_to_color(sev)

            bar = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top,
                Inches(0.1), Inches(2.3),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = sev_c
            bar.line.fill.background()

            _add_text(
                s, left + Inches(0.3), top - Inches(0.05),
                Inches(5.4), Inches(0.5),
                f"#{i+1}  {d.get('type', '')}",
                size=18, bold=True, color=DARK,
            )

            _add_text(
                s, left + Inches(0.3), top + Inches(0.45),
                Inches(2.5), Inches(0.3),
                f"SEVERITY  {sev}",
                size=10, bold=True, color=sev_c,
            )

            _add_text(
                s, left + Inches(0.3), top + Inches(0.8),
                Inches(5.4), Inches(0.3),
                f"위치: {d.get('location_desc', '-')}",
                size=11, color=DARK,
            )

            _add_text(
                s, left + Inches(0.3), top + Inches(1.1),
                Inches(5.4), Inches(0.3),
                f"크기: {d.get('size_estimate', '-')}",
                size=11, color=DARK,
            )

            desc = d.get("description", "")
            if len(desc) > 130:
                desc = desc[:130] + "..."
            _add_text(
                s, left + Inches(0.3), top + Inches(1.45),
                Inches(5.4), Inches(0.8),
                desc, size=11, color=DARK,
            )
    else:
        _add_text(
            s, Inches(0.7), Inches(3.5), Inches(12), Inches(1),
            "특이 결함이 검출되지 않았습니다",
            size=24, bold=True, color=SEV_LOW, align=PP_ALIGN.CENTER,
        )

    # ───── 슬라이드 4: 근본 원인 분석 ─────
    s = prs.slides.add_slide(blank)
    _fill_bg(s, WHITE, prs)

    _add_text(
        s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.3),
        "03  ·  ROOT CAUSE ANALYSIS",
        size=10, bold=True, color=NAVY,
    )
    _add_text(
        s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.6),
        "근본 원인 분석",
        size=28, bold=True, color=DARK,
    )
    _add_text(
        s, Inches(0.7), Inches(1.7), Inches(12), Inches(0.4),
        "웹 검색을 통해 자동 수집된 산업 표준 기반 분석",
        size=11, color=MUTED,
    )

    causes = research.get("root_cause_analysis", [])[:4]
    for i, c in enumerate(causes):
        row, col = divmod(i, 2)
        left = Inches(0.7 + col * 6.1)
        top = Inches(2.4 + row * 2.3)

        num = s.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.5), Inches(0.5))
        num.fill.solid()
        num.fill.fore_color.rgb = NAVY
        num.line.fill.background()
        tf = num.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        tf.text = str(i + 1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for r in p.runs:
            r.font.size = Pt(16)
            r.font.bold = True
            r.font.color.rgb = WHITE

        _add_text(
            s, left + Inches(0.7), top - Inches(0.02),
            Inches(5.0), Inches(0.45),
            c.get("category", ""),
            size=15, bold=True, color=NAVY,
        )

        details = c.get("details", "")
        if len(details) > 180:
            details = details[:180] + "..."
        _add_text(
            s, left + Inches(0.7), top + Inches(0.5),
            Inches(5.0), Inches(1.7),
            details, size=11, color=DARK,
        )

    # ───── 슬라이드 5: 권고 사항 ─────
    s = prs.slides.add_slide(blank)
    _fill_bg(s, WHITE, prs)

    _add_text(
        s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.3),
        "04  ·  RECOMMENDATIONS",
        size=10, bold=True, color=NAVY,
    )
    _add_text(
        s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.6),
        "권고 조치",
        size=28, bold=True, color=DARK,
    )

    actions = list(analysis.get("recommended_actions", []))
    actions += list(research.get("improvement_recommendations", []))
    actions = actions[:5]

    for i, a in enumerate(actions):
        top = Inches(1.95 + i * 1.0)
        _add_text(
            s, Inches(0.7), top - Inches(0.05),
            Inches(0.9), Inches(0.9),
            f"{i+1:02d}", size=32, bold=True, color=ICE,
        )
        if len(a) > 160:
            a = a[:160] + "..."
        _add_text(
            s, Inches(1.7), top + Inches(0.1),
            Inches(11.0), Inches(0.9),
            a, size=14, color=DARK,
        )

    # ───── 슬라이드 6: 참고 자료 (다크) ─────
    s = prs.slides.add_slide(blank)
    _fill_bg(s, NAVY, prs)

    _add_text(
        s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.3),
        "05  ·  REFERENCES",
        size=10, bold=True, color=ICE,
    )
    _add_text(
        s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.6),
        "참고 자료",
        size=28, bold=True, color=WHITE,
    )
    _add_text(
        s, Inches(0.7), Inches(1.85), Inches(12), Inches(0.4),
        "산업 표준 · 본 검사 보고서의 근거 자료",
        size=12, color=ICE,
    )

    _add_text(
        s, Inches(0.7), Inches(2.6), Inches(5.8), Inches(0.35),
        "INDUSTRY STANDARDS",
        size=11, bold=True, color=ICE,
    )
    standards = research.get("industry_standards", [])[:6]
    standards_text = "\n".join([f"·  {st}" for st in standards]) or "·  (해당 사항 없음)"
    _add_text(
        s, Inches(0.7), Inches(3.0), Inches(5.8), Inches(3.5),
        standards_text, size=10, color=ICE,
    )

    _add_text(
        s, Inches(6.9), Inches(2.6), Inches(5.8), Inches(0.35),
        "WEB SOURCES",
        size=11, bold=True, color=ICE,
    )
    sources = research.get("_sources", [])[:8]
    src_lines = []
    for i, src in enumerate(sources, 1):
        s_disp = src if len(src) <= 70 else src[:67] + "..."
        src_lines.append(f"[{i:02d}]  {s_disp}")
    src_text = "\n".join(src_lines) or "(검색 결과 없음)"
    _add_text(
        s, Inches(6.9), Inches(3.0), Inches(5.8), Inches(3.5),
        src_text, size=9, color=ICE,
    )

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# =============================================================================
# CLI
# =============================================================================
def main():
    p = argparse.ArgumentParser(
        description="Render inspection report PPT from JSON inputs."
    )
    p.add_argument("image", help="Path to inspection image (png/jpg/etc.)")
    p.add_argument("analysis", help="Path to analysis JSON")
    p.add_argument("research", help="Path to research JSON")
    p.add_argument("--product", default="PCB", help="Product type label for cover")
    p.add_argument("--out", required=True, help="Output .pptx path")
    p.add_argument(
        "--annotated-out",
        default=None,
        help="Optional path to save the annotated PNG alongside the PPT",
    )
    args = p.parse_args()

    image_path = Path(args.image)
    if not image_path.exists():
        raise SystemExit(f"Image not found: {image_path}")

    original_bytes = image_path.read_bytes()
    analysis = json.loads(Path(args.analysis).read_text(encoding="utf-8"))
    research = json.loads(Path(args.research).read_text(encoding="utf-8"))

    defects = analysis.get("defects", [])
    annotated_bytes = None
    if defects:
        annotated_bytes = draw_defect_overlay(original_bytes, defects)
        if args.annotated_out:
            Path(args.annotated_out).write_bytes(annotated_bytes)

    buf = generate_inspection_pptx(
        original_bytes, annotated_bytes, analysis, research, args.product
    )
    out_path = Path(args.out)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_bytes(buf.getvalue())

    print(f"Wrote {out_path}")
    if annotated_bytes and args.annotated_out:
        print(f"Wrote {args.annotated_out}")


if __name__ == "__main__":
    main()
